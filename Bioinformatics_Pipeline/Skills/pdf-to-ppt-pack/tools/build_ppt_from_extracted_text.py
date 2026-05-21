from __future__ import annotations

from datetime import datetime
from pathlib import Path
import argparse
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import random
import sys
import os
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn


FIGURE_RE = re.compile(r"Figure\s+(\d+)\.")
FIGURE_TITLE_RE = re.compile(r"Figure\s+(\d+)\.\s*([^\n]+)")
LABEL_PREFIX_RE = re.compile("^(Summary|[A-Z](?:\\s*[,–-]\\s*[A-Z])*)\\b")
LABEL_COLORS = [
    RGBColor(0x2C, 0x5F, 0x2D),
    RGBColor(0xF9, 0x61, 0x67),
    RGBColor(0x06, 0x5A, 0x82),
    RGBColor(0x99, 0x00, 0x11),
]
RANDOM_SEED_ENV = "PPT_RANDOM_SEED"

TRANSITION_PRESETS = ["fade", "push", "wipe"]
TITLE_ANIM_PRESETS = ["fadeIn", "slideUp", "slideRight", "scaleIn"]
TEXT_ANIM_PRESETS = ["fadeIn", "slideUp", "slideRight"]
IMAGE_ANIM_PRESETS = ["slideLeft", "fadeIn", "zoomIn"]


def build_animator():
    animations_dir = Path(__file__).resolve().parents[1] / "PPTX" / "animations"
    if str(animations_dir) not in sys.path:
        sys.path.insert(0, str(animations_dir))
    from animation import AnimationResolver, AnimationBuilder
    resolver = AnimationResolver(component_defaults={})
    return AnimationBuilder(resolver)


def apply_slide_transition(slide, transition_name):
    sld = slide._element
    if sld.find(qn("p:transition")) is not None:
        return
    transition = OxmlElement("p:transition")
    if transition_name == "push":
        transition.append(OxmlElement("p:push"))
    elif transition_name == "wipe":
        transition.append(OxmlElement("p:wipe"))
    else:
        transition.append(OxmlElement("p:fade"))
    timing = sld.find(qn("p:timing"))
    if timing is None:
        sld.append(transition)
    else:
        sld.insert(sld.index(timing), transition)


def animate_sequence(animator, slide, steps):
    if animator is None:
        return
    animator.reset()
    for shape, preset, step in steps:
        animator.queue_entrance(shape, preset, step)
    animator.apply_to_slide(slide)


def load_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_title(text: str) -> str:
    lines = [line.strip() for line in text.splitlines()]
    try:
        idx = lines.index("Article")
    except ValueError:
        return "Untitled Paper"
    title_lines = []
    for line in lines[idx + 1 : idx + 8]:
        if not line or line.lower().startswith("graphical abstract"):
            break
        if line.lower().startswith("highlights"):
            break
        title_lines.append(line)
    return " ".join(title_lines).strip() or "Untitled Paper"


def extract_journal_line(text: str) -> str:
    for line in text.splitlines():
        if "Cancer Cell" in line and re.search(r"\b\d{4}\b", line):
            return line.strip()
    for line in text.splitlines():
        if re.search(r"\b\d{4}\b", line) and "doi" in line.lower():
            return line.strip()
    return ""


def extract_section_lines(text: str, start: str, end_markers: list[str]) -> list[str]:
    lines = text.splitlines()
    start_idx = None
    for i, line in enumerate(lines):
        if line.strip().lower() == start.lower():
            start_idx = i + 1
            break
    if start_idx is None:
        return []
    collected = []
    for line in lines[start_idx:]:
        if not line.strip():
            if collected:
                continue
        if any(line.strip().lower().startswith(m.lower()) for m in end_markers):
            break
        collected.append(line.strip())
    return [l for l in collected if l]


def split_sentences(text: str, max_sentences: int = 4) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    parts = re.split(r"(?<=[.!?])\s+", text)
    return [p.strip() for p in parts if p.strip()][:max_sentences]


def extract_figures(text: str) -> dict[str, str]:
    matches = list(FIGURE_RE.finditer(text))
    figures: dict[str, str] = {}
    for i, match in enumerate(matches):
        fig_no = match.group(1)
        start = match.start()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        legend = text[start:end]
        legend = legend.split("OPEN ACCESS")[0]
        legend = legend.replace("\n", " ")
        legend = re.sub(r"\s+", " ", legend).strip()
        figures[f"Figure_{fig_no}"] = legend
    return figures


def extract_figure_titles(text: str) -> dict[str, str]:
    raw = text.replace("-\n", "").replace("\n", " ")
    raw = re.sub(r"\s+", " ", raw)
    titles = {}
    for match in FIGURE_TITLE_RE.finditer(raw):
        fig_no = match.group(1)
        title = match.group(2).strip().rstrip(".")
        titles[f"Figure_{fig_no}"] = title
    return titles


def load_figure_titles_zh(fig_dir: Path) -> dict[str, str]:
    zh_path = fig_dir / "figure_titles_zh.txt"
    if not zh_path.exists():
        return {}
    zh_titles = {}
    for line in zh_path.read_text(encoding="utf-8", errors="ignore").splitlines():
        if not line.strip():
            continue
        if "\t" not in line:
            continue
        key, title = line.split("\t", 1)
        zh_titles[key.strip()] = title.strip()
    return zh_titles


def build_summary_line(bullets: list[str]) -> str:
    for line in bullets:
        cleaned = re.sub(r"^#+\s*", "", line).strip()
        cleaned = LABEL_PREFIX_RE.sub("", cleaned).lstrip(" ,–-")
        if not cleaned:
            continue
        if "." in cleaned:
            snippet = cleaned.split(".", 1)[0]
        else:
            snippet = cleaned[:80]
        return f"summary：{snippet}"
    return "Summary: The results of this figure support this conclusion."


def add_title_slide(
    prs: Presentation,
    title: str,
    journal_date: str,
    animator=None,
    transition_name="fade",
    anim_presets=None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(2.0825), Inches(1.796), Inches(9.1683), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    if journal_date:
        tb2 = slide.shapes.add_textbox(Inches(4.8673), Inches(4.5049), Inches(6.0), Inches(0.6))
        tf2 = tb2.text_frame
        tf2.paragraphs[0].text = journal_date
        tf2.paragraphs[0].font.size = Pt(18)

    tb3 = slide.shapes.add_textbox(Inches(4.8673), Inches(5.2327), Inches(6.0), Inches(0.6))
    tf3 = tb3.text_frame
    tf3.paragraphs[0].text = "Presenter: XXX"
    tf3.paragraphs[0].font.size = Pt(16)

    tb4 = slide.shapes.add_textbox(Inches(4.8673), Inches(5.9049), Inches(6.0), Inches(0.6))
    tf4 = tb4.text_frame
    tf4.paragraphs[0].text = f"Report date: {datetime.now().strftime('%Y.%m.%d')}"
    tf4.paragraphs[0].font.size = Pt(16)

    presets = anim_presets or {}
    apply_slide_transition(slide, transition_name)
    steps = [(tb, presets.get("title", "fadeIn"), 1)]
    if journal_date:
        steps.append((tb2, presets.get("text", "fadeIn"), 2))
    steps.extend([
        (tb3, presets.get("text", "fadeIn"), 3),
        (tb4, presets.get("text", "fadeIn"), 3),
    ])
    animate_sequence(animator, slide, steps)


def add_overview_slide(
    prs: Presentation,
    background_bullets: list[str],
    findings_bullets: list[str],
    ga_image_path: Path | None,
    animator=None,
    transition_name="fade",
    anim_presets=None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb_left = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(5.8), Inches(6.2))
    tf_left = tb_left.text_frame
    p = tf_left.paragraphs[0]
    p.text = "Background"
    p.font.size = Pt(20)
    p.font.bold = True
    for item in background_bullets:
        para = tf_left.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    tb_mid = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(5.8), Inches(3.4))
    tf_mid = tb_mid.text_frame
    p2 = tf_mid.paragraphs[0]
    p2.text = "Key Findings"
    p2.font.size = Pt(20)
    p2.font.bold = True
    for item in findings_bullets:
        para = tf_mid.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    if ga_image_path and ga_image_path.exists():
        ga_shape = slide.shapes.add_picture(str(ga_image_path), Inches(7.0), Inches(0.8), Inches(5.8), Inches(5.6))
    else:
        ga_shape = None

    presets = anim_presets or {}
    apply_slide_transition(slide, transition_name)
    steps = [
        (tb_left, presets.get("title", "fadeIn"), 1),
        (tb_mid, presets.get("text", "fadeIn"), 2),
    ]
    if ga_shape is not None:
        steps.append((ga_shape, presets.get("image", "fadeIn"), 3))
    animate_sequence(animator, slide, steps)


def add_result_slide(
    prs: Presentation,
    title: str,
    image_path: Path | None,
    bullets: list[str],
    label_color: RGBColor,
    animator=None,
    transition_name="fade",
    anim_presets=None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    tb2 = slide.shapes.add_textbox(Inches(0.7749), Inches(1.5102), Inches(4.7498), Inches(2.861))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, line in enumerate(bullets):
        para = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        cleaned = re.sub(r"^#+\s*", "", line).strip()
        para.text = ""
        match = LABEL_PREFIX_RE.match(cleaned)
        if match:
            prefix = match.group(1)
            rest = cleaned[len(prefix):].lstrip(" ,–-")
            run_prefix = para.add_run()
            run_prefix.text = prefix
            run_prefix.font.bold = True
            run_prefix.font.size = Pt(14)
            run_prefix.font.color.rgb = label_color
            run_rest = para.add_run()
            run_rest.text = f" {rest}".rstrip()
            run_rest.font.size = Pt(14)
        else:
            run = para.add_run()
            run.text = cleaned
            run.font.size = Pt(14)
        para.space_after = Pt(4)

    if image_path and image_path.exists():
        img_shape = slide.shapes.add_picture(str(image_path), Inches(7.0), Inches(1.0), Inches(5.8), Inches(5.8))
    else:
        img_shape = None

    presets = anim_presets or {}
    apply_slide_transition(slide, transition_name)
    steps = [
        (tb, presets.get("title", "fadeIn"), 1),
        (tb2, presets.get("text", "fadeIn"), 2),
    ]
    if img_shape is not None:
        steps.append((img_shape, presets.get("image", "fadeIn"), 3))
    animate_sequence(animator, slide, steps)


def add_summary_slide(
    prs: Presentation,
    summary_bullets: list[str],
    notes_bullets: list[str],
    label_color: RGBColor,
    animator=None,
    transition_name="fade",
    anim_presets=None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box_fill = RGBColor(0xF2, 0xF2, 0xF2)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(5.8), Inches(3.0))
    tb.fill.solid()
    tb.fill.fore_color.rgb = box_fill
    tb.line.color.rgb = box_fill
    tb.line.width = Pt(1)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Summary"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = label_color
    for item in summary_bullets:
        para = tf.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    tb2 = slide.shapes.add_textbox(Inches(6.7), Inches(3.9), Inches(6.0), Inches(3.0))
    tb2.fill.solid()
    tb2.fill.fore_color.rgb = box_fill
    tb2.line.color.rgb = box_fill
    tb2.line.width = Pt(1)
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Notes"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = label_color
    for item in notes_bullets:
        para = tf2.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    presets = anim_presets or {}
    apply_slide_transition(slide, transition_name)
    animate_sequence(
        animator,
        slide,
        [
            (tb, presets.get("title", "fadeIn"), 1),
            (tb2, presets.get("text", "fadeIn"), 2),
        ],
    )


def main() -> None:
    parser = argparse.ArgumentParser(description="Build PPT from extracted_text.txt and Figure_*.jpg files.")
    parser.add_argument("--text", required=True, help="Path to extracted_text.txt")
    parser.add_argument("--fig-dir", required=True, help="Directory containing Figure_*.jpg")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    args = parser.parse_args()

    text_path = Path(args.text)
    fig_dir = Path(args.fig_dir)
    output_path = Path(args.output)

    seed = os.getenv(RANDOM_SEED_ENV)
    if seed:
        try:
            random.seed(int(seed))
        except ValueError:
            random.seed(seed)

    text = load_text(text_path)
    title = extract_title(text)
    journal_date = extract_journal_line(text)

    highlights = extract_section_lines(text, "Highlights", ["Authors", "In brief"])
    in_brief_lines = extract_section_lines(text, "In brief", ["Authors", "SUMMARY"])
    summary_lines = extract_section_lines(text, "SUMMARY", ["INTRODUCTION"])

    background = split_sentences(" ".join(in_brief_lines), max_sentences=3) or [
        "Background not found in extracted text.",
    ]
    findings = highlights or split_sentences(" ".join(summary_lines), max_sentences=3)
    if not findings:
        findings = ["Key findings not found in extracted text."]

    figures = extract_figures(text)
    figure_titles = extract_figure_titles(text)
    figure_titles_zh = load_figure_titles_zh(fig_dir)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    animator = build_animator()
    transition_name = random.choice(TRANSITION_PRESETS)
    anim_presets = {
        "title": random.choice(TITLE_ANIM_PRESETS),
        "text": random.choice(TEXT_ANIM_PRESETS),
        "image": random.choice(IMAGE_ANIM_PRESETS),
    }
    add_title_slide(prs, title, journal_date, animator, transition_name, anim_presets)

    ga_path = fig_dir / "Graphical abstract.jpg"
    add_overview_slide(
        prs,
        background,
        findings,
        ga_path if ga_path.exists() else None,
        animator,
        transition_name,
        anim_presets,
    )

    fig_files = []
    for fig_path in fig_dir.glob("Figure_*.jpg"):
        match = re.search(r"(\\d+)", fig_path.stem)
        if not match:
            continue
        fig_files.append((int(match.group(1)), fig_path))
    fig_files.sort(key=lambda item: item[0])
    label_color = random.choice(LABEL_COLORS)
    for _num, fig_path in fig_files:
        fig_label = fig_path.stem
        legend = figures.get(fig_label, "")
        bullets = split_sentences(legend, max_sentences=5)
        if not bullets:
            bullets = [f"{fig_label} legend not found in extracted text."]
        bullets.append(build_summary_line(bullets))
        fig_title = figure_titles_zh.get(fig_label) or figure_titles.get(fig_label)
        if fig_title:
            slide_title = f"{fig_label.replace('_', ' ')}: {fig_title}"
        else:
            slide_title = fig_label.replace("_", " ")
        add_result_slide(
            prs,
            slide_title,
            fig_path,
            bullets,
            label_color,
            animator,
            transition_name,
            anim_presets,
        )

    summary_bullets = split_sentences(" ".join(summary_lines), max_sentences=4)
    if not summary_bullets:
        summary_bullets = ["Summary not found in extracted text."]
    notes_bullets = ["Verify figure interpretation and refine slide text as needed."]
    add_summary_slide(prs, summary_bullets, notes_bullets, label_color, animator, transition_name, anim_presets)

    prs.save(output_path)


if __name__ == "__main__":
    main()
