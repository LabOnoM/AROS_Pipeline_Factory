from datetime import datetime
from pathlib import Path
import argparse
import re

from pptx import Presentation
from pptx.util import Inches, Pt


def parse_md_pages(md_text: str) -> dict[int, str]:
    pages = {}
    parts = md_text.split("## Page ")
    for part in parts[1:]:
        header, _, body = part.partition("\n")
        try:
            page_num = int(header.strip()[:2])
        except ValueError:
            continue
        pages[page_num] = body.strip()
    return pages


def extract_title_and_journal(page1_text: str) -> tuple[str, str]:
    title = "Hepatic ER stress suppresses adipose browning through ATF4-CIRP-ANGPTL3 cascade"
    journal = "Cell Reports, 2022 Sep 29"
    m_title = re.search(r"Article\s+(.+?)\s+Graphical abstract", page1_text, re.IGNORECASE)
    if m_title:
        title = m_title.group(1).strip()
    m_journal = re.search(r"(Cell Reports[^\\n]*?\\d{4})", page1_text)
    if m_journal:
        journal = m_journal.group(1).strip()
    return title, journal


def extract_highlights(page1_text: str) -> list[str]:
    if "Highlights" not in page1_text:
        return []
    _, _, tail = page1_text.partition("Highlights")
    # Split by " d " markers used in the PDF text.
    items = [s.strip() for s in re.split(r"\\bd\\b", tail) if s.strip()]
    highlights = []
    for item in items:
        # Stop if we hit another section marker
        if item.lower().startswith("lv et al") or "doi" in item.lower():
            break
        highlights.append(item.replace("  ", " ").strip())
    return highlights[:4]


def map_figure_pages(pages: dict[int, str]) -> dict[int, int]:
    fig_pages = {}
    for page_num, body in pages.items():
        for m in re.finditer(r"Figure\\s+(\\d+)\\.", body):
            fig_no = int(m.group(1))
            if fig_no not in fig_pages:
                fig_pages[fig_no] = page_num
    return fig_pages


def add_title_slide(prs, title, journal_date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12.0), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.6))
    tf2 = tb2.text_frame
    tf2.paragraphs[0].text = journal_date
    tf2.paragraphs[0].font.size = Pt(18)

    tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(6.0), Inches(0.6))
    tf3 = tb3.text_frame
    tf3.paragraphs[0].text = "Reporter:XXX"
    tf3.paragraphs[0].font.size = Pt(16)

    tb4 = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(6.0), Inches(0.6))
    tf4 = tb4.text_frame
    tf4.paragraphs[0].text = f"Reporting date：{datetime.now().strftime('%Y.%m.%d')}"
    tf4.paragraphs[0].font.size = Pt(16)


def add_overview_slide(prs, background_bullets, findings_bullets, ga_image_path: Path | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb_left = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(5.8), Inches(2.6))
    tf_left = tb_left.text_frame
    p = tf_left.paragraphs[0]
    p.text = "Research background"
    p.font.size = Pt(20)
    p.font.bold = True
    for item in background_bullets:
        para = tf_left.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    tb_mid = slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.6))
    tf_mid = tb_mid.text_frame
    p2 = tf_mid.paragraphs[0]
    p2.text = "Core findings"
    p2.font.size = Pt(20)
    p2.font.bold = True
    for item in findings_bullets:
        para = tf_mid.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    if ga_image_path and ga_image_path.exists():
        slide.shapes.add_picture(str(ga_image_path), Inches(7.0), Inches(0.8), Inches(5.8), Inches(5.6))


def add_result_slide(prs, title, image_path, paragraphs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True

    tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(6.2), Inches(6.0))
    tf2 = tb2.text_frame
    for idx, line in enumerate(paragraphs):
        para = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        para.text = line
        para.font.size = Pt(14)
        para.space_after = Pt(4)

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(7.0), Inches(1.0), Inches(5.8), Inches(5.8))


def add_summary_slide(prs, strengths, limits):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(6.0), Inches(6.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Highlights"
    p.font.size = Pt(20)
    p.font.bold = True
    for item in strengths:
        para = tf.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    tb2 = slide.shapes.add_textbox(Inches(7.0), Inches(0.6), Inches(6.0), Inches(6.2))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "inherent limitations"
    p2.font.size = Pt(20)
    p2.font.bold = True
    for item in limits:
        para = tf2.add_paragraph()
        para.text = item
        para.font.size = Pt(16)


def main():
    parser = argparse.ArgumentParser(description="Build PPT from pdf-extract Markdown and page images.")
    parser.add_argument("--md", required=True, help="Path to extracted Markdown.")
    parser.add_argument("--image-dir", required=True, help="Directory with page-XX.png images.")
    parser.add_argument("--output", required=True, help="Output PPTX path.")
    args = parser.parse_args()

    md_path = Path(args.md)
    image_dir = Path(args.image_dir)
    output_path = Path(args.output)

    md_text = md_path.read_text(encoding="utf-8", errors="ignore")
    pages = parse_md_pages(md_text)
    title, journal = extract_title_and_journal(pages.get(1, ""))
    highlights = extract_highlights(pages.get(1, ""))

    background = [
        "Hepatic ER stress is closely associated with obesity/fatty liver disease",
        "Liver-derived factors may regulate fat browning through cross-organ communication",
    ]
    findings = highlights if highlights else [
        "ATF4 upregulates CIRP and promotes ANGPTL3 secretion",
        "ANGPTL3 inhibits fat browning through integrin αvβ3-JNK",
    ]

    fig_pages = map_figure_pages(pages)

    results = {
        "Figure_1": [
            "# A, B Comparing the expression of CIRP in livers of HFD and RD mice, RT-qPCR and Western blotting showed that liver CIRP mRNA and protein were up-regulated in the HFD group, but there were no changes in adipose tissue.",
            "# C, D Hepatic CIRP mRNA and protein were significantly increased in ob/ob mice.",
            "# E, F The liver CIRP mRNA and protein of obese people are higher than those of lean people.",
            "# G CIRP mRNA was significantly positively correlated with BMI.",
        ],
        "Figure_2": [
            "# A Liver AAV-CIRPi knocks down CIRP, and liver CIRP protein decreases.",
            "# B There was no significant difference in food intake.",
            "# C, D Body weight gain is reduced and scWAT weight is reduced.",
            "# E scWAT adipocytes become smaller and UCP1 staining is enhanced.",
            "# F The thermogenic gene Ucp1/Pgc1a/Cidea is up-regulated.",
            "# G scWAT UCP1 protein is elevated.",
            "# H VO2 increases and energy consumption increases.",
        ],
        "Figure_3": [
            "# A Ad-CIRP increases liver CIRP protein.",
            "# B VO2 decreases under cold exposure.",
            "# C scWAT thermogenic gene is downregulated.",
            "# D, E scWAT UCP1 protein decreased and immunohistochemistry was weakened.",
            "# F scWAT OCR decreased.",
            "# G The VO2 difference disappears in the UCP1 KO background.",
        ],
        "Figure_4": [
            "# A ANGPTL3 is the most significantly up-regulated liver-derived factor induced by Ad-CIRP.",
            "# B, C Ad-CIRP increases ANGPTL3 secretion, while CIRP KO decreases secretion.",
            "# D, E Changes in liver ANGPTL3 protein with CIRP upregulation or knockout.",
            "# F-I CIRP stabilizes ANGPTL3 mRNA and enhances reporter activity by binding to its 3'UTR.",
            "# K-M The liver ANGPTL3 protein is increased in obesity models and obese people.",
        ],
        "Figure_5": [
            "# A Ad-ANGPTL3 reduces VO2.",
            "# B-D scWAT thermogenic gene and UCP1 protein were down-regulated, and histology showed reduced browning.",
            "# E scWAT OCR decreased.",
            "# F, G Recombinant ANGPTL3 dose-dependently inhibits thermogenic genes and reduces UCP1 protein.",
            "# H-K ANGPTL3 KO weakens the inhibitory effect of CIRP on VO2 and UCP1.",
        ],
        "Figure_6": [
            "# A reANGPTL3 can still further inhibit thermogenic genes after orlistat.",
            "# B Integrin αvβ3 is upregulated during differentiation.",
            "# C cilengitide inhibits reANGPTL3-induced FAK/JNK phosphorylation.",
            "# D, E cilengitide or JNK inhibitor relieves the inhibitory effect of ANGPTL3.",
            "# F-I ob/ob mice were injected with cilengitide, VO2 and UCP1 were up-regulated, and adipocytes became smaller.",
        ],
        "Figure_7": [
            "# A-C ATF4-activated conditioned medium inhibits thermogenic genes and UCP1.",
            "# D, E Hepatic ATF4 is upregulated and VO2 is reduced.",
            "# F-H scWAT UCP1 and OCR down.",
            "# I BAT maximum OCR decreases.",
            "# J, K ATF4 upregulation induces the expression of CIRP and ANGPTL3.",
        ],
    }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, title, journal)

    ga_path = image_dir / "page-01.png"
    add_overview_slide(prs, background, findings, ga_path)

    for fig_no in range(1, 8):
        fig_key = f"Figure_{fig_no}"
        page_num = fig_pages.get(fig_no)
        image_path = image_dir / f"page-{page_num:02d}.png" if page_num else None
        add_result_slide(prs, fig_key, image_path, results[fig_key])

    strengths = [
        "Propose the ATF4-CIRP-ANGPTL3 cascade as a fat browning inhibitory pathway",
        "Multi-level experiments verify cross-organ regulation of liver-derived factors",
        "Integrin αvβ3/JNK as a potential intervention target",
    ]
    limits = [
        "Human evidence is limited and needs to be verified with larger samples",
        "The upstream triggering mechanism still needs to be supplemented",
        "The safety and translatability of long-term intervention need to be evaluated",
    ]
    add_summary_slide(prs, strengths, limits)

    prs.save(str(output_path))


if __name__ == "__main__":
    main()
