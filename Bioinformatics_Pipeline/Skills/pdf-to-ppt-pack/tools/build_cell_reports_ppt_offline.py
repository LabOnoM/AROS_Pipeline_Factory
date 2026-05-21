from datetime import datetime
from pathlib import Path
import argparse
import re
import random
import sys
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn


DEFAULT_BASE_DIR = Path(r"d:\SKILL\project\cell-reports-36170814_offline")
DEFAULT_OUTPUT_PATH = Path(r"d:\SKILL\project\cell-reports-36170814.pptx")

FIGURE_TITLE_RE = re.compile(r"Figure\s+(\d+)\.\s*([^\n]+)")
LABEL_PREFIX_RE = re.compile("^(Summary|[A-Z](?:\\s*[,–-]\\s*[A-Z])*)\\b")
LABEL_COLORS = [
    RGBColor(0x2C, 0x5F, 0x2D),
    RGBColor(0xF9, 0x61, 0x67),
    RGBColor(0x06, 0x5A, 0x82),
    RGBColor(0x99, 0x00, 0x11),
]
RANDOM_SEED_ENV = "PPT_RANDOM_SEED"

TRANSITION_PRESETS = ["fade", "push", "wipe"]
TITLE_ANIM_PRESETS = ["fadeIn", "slideUp", "slideRight", "scaleIn"]
TEXT_ANIM_PRESETS = ["fadeIn", "slideUp", "slideRight"]
IMAGE_ANIM_PRESETS = ["slideLeft", "fadeIn", "zoomIn"]


def build_animator():
    animations_dir = Path(__file__).resolve().parents[1] / "PPTX" / "animations"
    if not animations_dir.exists():
        animations_dir = Path(__file__).resolve().parents[1] / "ppt" / "framework" / "converter"
    if not animations_dir.exists():
        return None
    if str(animations_dir) not in sys.path:
        sys.path.insert(0, str(animations_dir))
    try:
        from animation import AnimationResolver, AnimationBuilder
    except Exception:
        return None
    resolver = AnimationResolver(component_defaults={})
    return AnimationBuilder(resolver)


def apply_slide_transition(slide, transition_name):
    sld = slide._element
    if sld.find(qn("p:transition")) is not None:
        return
    transition = OxmlElement("p:transition")
    if transition_name == "push":
        transition.append(OxmlElement("p:push"))
    elif transition_name == "wipe":
        transition.append(OxmlElement("p:wipe"))
    else:
        transition.append(OxmlElement("p:fade"))
    timing = sld.find(qn("p:timing"))
    if timing is None:
        sld.append(transition)
    else:
        sld.insert(sld.index(timing), transition)


def animate_sequence(animator, slide, steps):
    if animator is None:
        return
    animator.reset()
    for shape, preset, step in steps:
        animator.queue_entrance(shape, preset, step)
    animator.apply_to_slide(slide)


def load_figure_titles(base_dir: Path) -> dict[str, str]:
    text_path = base_dir / "extracted_text.txt"
    if not text_path.exists():
        return {}
    raw = text_path.read_text(encoding="utf-8", errors="ignore")
    raw = raw.replace("-\n", "").replace("\n", " ")
    raw = re.sub(r"\s+", " ", raw)
    titles = {}
    for match in FIGURE_TITLE_RE.finditer(raw):
        fig_no = match.group(1)
        title = match.group(2).strip().rstrip(".")
        titles[f"Figure_{fig_no}"] = title
    return titles


def load_figure_titles_zh(base_dir: Path) -> dict[str, str]:
    zh_path = base_dir / "figure_titles_zh.txt"
    if not zh_path.exists():
        return {}
    zh_titles = {}
    for line in zh_path.read_text(encoding="utf-8", errors="ignore").splitlines():
        if not line.strip():
            continue
        if "\t" not in line:
            continue
        key, title = line.split("\t", 1)
        zh_titles[key.strip()] = title.strip()
    return zh_titles


def normalize_figure_title(title: str, max_len: int = 120) -> str:
    if not title:
        return title
    title = title.strip()
    for marker in [" (A", " (a", " (", " (A)", " (a)"]:
        if marker in title:
            title = title.split(marker, 1)[0].strip()
            break
    if ". " in title:
        title = title.split(". ", 1)[0].strip()
    if len(title) > max_len:
        title = title[:max_len].rstrip()
    return title


def build_summary_line(paragraphs: list[str]) -> str:
    for line in paragraphs:
        cleaned = re.sub(r"^#+\s*", "", line).strip()
        cleaned = LABEL_PREFIX_RE.sub("", cleaned).lstrip(" ,–-")
        if not cleaned:
            continue
        if "。" in cleaned:
            snippet = cleaned.split("。", 1)[0]
        else:
            snippet = cleaned[:60]
        return f"summary：{snippet}"
    return "Summary: The results of this figure support this conclusion."


def add_title_slide(prs, title, journal_date, animator=None, transition_name="fade", anim_presets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    tb = slide.shapes.add_textbox(Inches(2.0825), Inches(1.796), Inches(9.1683), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Journal/date
    tb2 = slide.shapes.add_textbox(Inches(4.8673), Inches(4.5049), Inches(6.0), Inches(0.6))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = journal_date
    p2.font.size = Pt(18)

    # Presenter and report date
    tb3 = slide.shapes.add_textbox(Inches(4.8673), Inches(5.2327), Inches(6.0), Inches(0.6))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"Reporter：XXX"
    p3.font.size = Pt(16)

    tb4 = slide.shapes.add_textbox(Inches(4.8673), Inches(5.9049), Inches(6.0), Inches(0.6))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = f"Reporting date：{datetime.now().strftime('%Y.%m.%d')}"
    p4.font.size = Pt(16)

    presets = anim_presets or {}
    apply_slide_transition(slide, transition_name)
    animate_sequence(
        animator,
        slide,
        [
            (tb, presets.get("title", "fadeIn"), 1),
            (tb2, presets.get("text", "fadeIn"), 2),
            (tb3, presets.get("text", "fadeIn"), 3),
            (tb4, presets.get("text", "fadeIn"), 3),
        ],
    )


def add_overview_slide(prs, background_bullets, findings_bullets, base_dir, animator=None, transition_name="fade", anim_presets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Left column background
    tb_left = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(5.8), Inches(6.2))
    tf_left = tb_left.text_frame
    p = tf_left.paragraphs[0]
    p.text = "Research background"
    p.font.size = Pt(20)
    p.font.bold = True
    for item in background_bullets:
        para = tf_left.add_paragraph()
        para.text = item
        para.level = 0
        para.font.size = Pt(16)

    # Left column findings
    tb_mid = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(5.8), Inches(3.4))
    tf_mid = tb_mid.text_frame
    p2 = tf_mid.paragraphs[0]
    p2.text = "Core findings"
    p2.font.size = Pt(20)
    p2.font.bold = True
    for item in findings_bullets:
        para = tf_mid.add_paragraph()
        para.text = item
        para.level = 0
        para.font.size = Pt(16)

    # Right image
    ga_path = base_dir / "Graphical abstract.jpg"
    if ga_path.exists():
        ga_shape = slide.shapes.add_picture(str(ga_path), Inches(7.0), Inches(0.8), Inches(5.8), Inches(5.6))
    else:
        ga_shape = None

    presets = anim_presets or {}
    apply_slide_transition(slide, transition_name)
    steps = [
        (tb_left, presets.get("title", "fadeIn"), 1),
        (tb_mid, presets.get("text", "fadeIn"), 2),
    ]
    if ga_shape is not None:
        steps.append((ga_shape, presets.get("image", "fadeIn"), 3))
    animate_sequence(animator, slide, steps)


def add_result_slide(prs, title, image_path, paragraphs, label_color, animator=None, transition_name="fade", anim_presets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Text box (left)
    tb2 = slide.shapes.add_textbox(Inches(0.7749), Inches(1.5102), Inches(4.7498), Inches(2.861))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, line in enumerate(paragraphs):
        cleaned = re.sub(r"^#+\s*", "", line).strip()
        if idx == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = ""
        match = LABEL_PREFIX_RE.match(cleaned)
        if match:
            prefix = match.group(1)
            rest = cleaned[len(prefix):].lstrip(" ,–-")
            run_prefix = p2.add_run()
            run_prefix.text = prefix
            run_prefix.font.bold = True
            run_prefix.font.size = Pt(14)
            run_prefix.font.color.rgb = label_color
            run_rest = p2.add_run()
            run_rest.text = f" {rest}".rstrip()
            run_rest.font.size = Pt(14)
        else:
            run = p2.add_run()
            run.text = cleaned
            run.font.size = Pt(14)
        p2.space_after = Pt(4)

    # Image (right)
    if image_path and Path(image_path).exists():
        img_shape = slide.shapes.add_picture(str(image_path), Inches(7.0), Inches(1.0), Inches(5.8), Inches(5.8))
    else:
        img_shape = None

    presets = anim_presets or {}
    apply_slide_transition(slide, transition_name)
    steps = [
        (tb, presets.get("title", "fadeIn"), 1),
        (tb2, presets.get("text", "fadeIn"), 2),
    ]
    if img_shape is not None:
        steps.append((img_shape, presets.get("image", "fadeIn"), 3))
    animate_sequence(animator, slide, steps)


def add_summary_slide(prs, strengths, limits, label_color, animator=None, transition_name="fade", anim_presets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(5.8), Inches(3.0))
    box_fill = RGBColor(0xF2, 0xF2, 0xF2)
    tb.fill.solid()
    tb.fill.fore_color.rgb = box_fill
    tb.line.color.rgb = box_fill
    tb.line.width = Pt(1)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Highlights"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = label_color
    for item in strengths:
        para = tf.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    tb2 = slide.shapes.add_textbox(Inches(6.7), Inches(3.9), Inches(6.0), Inches(3.0))
    tb2.fill.solid()
    tb2.fill.fore_color.rgb = box_fill
    tb2.line.color.rgb = box_fill
    tb2.line.width = Pt(1)
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "inherent limitations"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = label_color
    for item in limits:
        para = tf2.add_paragraph()
        para.text = item
        para.font.size = Pt(16)

    presets = anim_presets or {}
    apply_slide_transition(slide, transition_name)
    animate_sequence(
        animator,
        slide,
        [
            (tb, presets.get("title", "fadeIn"), 1),
            (tb2, presets.get("text", "fadeIn"), 2),
        ],
    )


def main():
    parser = argparse.ArgumentParser(description="Build PPT for cell-reports-36170814 (offline parse).")
    parser.add_argument("--base-dir", default=str(DEFAULT_BASE_DIR), help="Folder containing Figure_*.jpg")
    parser.add_argument("--output", default=str(DEFAULT_OUTPUT_PATH), help="Output PPTX path")
    args = parser.parse_args()

    base_dir = Path(args.base_dir)
    output_path = Path(args.output)
    figure_titles = load_figure_titles(base_dir)
    figure_titles_zh = load_figure_titles_zh(base_dir)

    seed = os.getenv(RANDOM_SEED_ENV)
    if seed:
        try:
            random.seed(int(seed))
        except ValueError:
            random.seed(seed)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = "Hepatic ER stress suppresses adipose browning through ATF4-CIRP-ANGPTL3 cascade"
    journal_date = "Cell Reports, 2022 Sep 29"

    animator = build_animator()
    transition_name = random.choice(TRANSITION_PRESETS)
    anim_presets = {
        "title": random.choice(TITLE_ANIM_PRESETS),
        "text": random.choice(TEXT_ANIM_PRESETS),
        "image": random.choice(IMAGE_ANIM_PRESETS),
    }
    add_title_slide(prs, title, journal_date, animator, transition_name, anim_presets)

    # Highlights from PDF first page
    background = [
        "Obesity-related hepatic stress may inhibit fat browning and energy expenditure",
        "Liver-secreted factors may play key roles in cross-organ metabolism regulation",
    ]
    findings = [
        "CIRP upregulation is associated with obesity and inhibits fat browning",
        "The ATF4-CIRP-ANGPTL3 cascade pathway drives this inhibition",
    ]
    add_overview_slide(prs, background, findings, base_dir, animator, transition_name, anim_presets)

    results = {
        "Figure_1": [
            "# A, B To compare the expression of CIRP in HFD and RD mice, RT-qPCR and Western blotting were used to detect liver and adipose tissue. The results showed that liver CIRP mRNA and protein were up-regulated in the HFD group, while scWAT/epiWAT/BAT had no significant changes.",
            "# C, D CIRP expression was detected in ob/ob mice. The results showed that liver CIRP mRNA and protein were significantly increased compared with WT, but no significant changes were found in adipose tissue.",
            "# E, F Comparing liver CIRP levels between lean people and obese people, the results show that liver CIRP mRNA and protein are significantly up-regulated in obese people.",
            "# G performed correlation analysis between BMI and liver CIRP mRNA, and the results showed that there was a significant positive correlation between the two.",
        ],
        "Figure_2": [
            "# A AAV-CIRPi was used to specifically knock down liver CIRP. The results showed that liver CIRP protein was down-regulated and scWAT/BAT remained unchanged.",
            "# B Comparison of food intake showed no significant difference between the AAV-CIRPi group and the control group.",
            "# C, D The AAV-CIRPi group had reduced body weight gain and decreased scWAT weight, suggesting that fat accumulation was inhibited.",
            "# E scWAT adipocyte volume decreased and UCP1 immunostaining increased, indicating enhanced browning; BAT showed no significant changes.",
            "# F The expression of thermogenic genes such as Ucp1, Pgc1a, Cidea is up-regulated in scWAT.",
            "# G UCP1 protein levels are elevated in scWAT.",
            "# H Indirect calorimetry shows increased VO2 and increased energy consumption.",
        ],
        "Figure_3": [
            "# A Hepatic Ad-CIRP overexpression significantly increases CIRP protein in the liver, while scWAT/BAT remains unchanged.",
            "# B VO2 decreased in the Ad-CIRP group under cold exposure conditions.",
            "# C scWAT thermogenic gene expression is down-regulated.",
            "# D, E UCP1 protein levels decreased in scWAT, and immunohistochemistry showed that UCP1 staining was weakened.",
            "# F Seahorse detection shows reduced scWAT base and maximum OCR.",
            "# G In UCP1 KO mice, the effect of Ad-CIRP on VO2 disappeared, suggesting that the effect is dependent on UCP1.",
        ],
        "Figure_4": [
            "# A Screening of various liver-derived factor mRNAs, ANGPTL3 was the most significantly up-regulated in the Ad-CIRP group.",
            "# B, C The secretion of ANGPTL3 increased after Ad-CIRP infected primary hepatocytes, while the secretion of CIRP KO cells decreased.",
            "# D, E ANGPTL3 protein in the liver increased in the Ad-CIRP group and decreased in the CIRP KO group.",
            "# F shows potential sequences for CIRP binding ANGPTL3 3'UTR.",
            "# G, H mRNA half-life experiments showed that CIRP overexpression extended the half-life of ANGPTL3 mRNA, while CIRP knockdown shortened the half-life.",
            "# I, J 3'UTR luciferase experiment shows that CIRP enhances the activity of WT 3'UTR, and the effect disappears after mutating the binding site.",
            "# K-M HFD, ob/ob mice and obese people have up-regulated liver ANGPTL3 protein.",
        ],
        "Figure_5": [
            "# A Ad-ANGPTL3 overexpression reduces VO2.",
            "# B-D scWAT thermogenic gene and UCP1 protein were down-regulated, and histology showed reduced browning.",
            "# E scWAT base and maximum OCR reduced.",
            "# F, G Recombinant ANGPTL3 treatment of primary adipocytes dose-dependently inhibits thermogenic genes and reduces UCP1 protein.",
            "# H-K In the ANGPTL3 KO background, the inhibitory effect of Ad-CIRP on VO2 and scWAT thermogenic gene/UCP1 is weakened or disappeared.",
        ],
        "Figure_6": [
            "# A Orlistat inhibits LPL and reduces thermogenic genes, but reANGPTL3 can further inhibit it, suggesting the existence of a non-LPL pathway.",
            "# B Itga5/Itgb3 is up-regulated during beige fat differentiation, suggesting the involvement of integrin αvβ3.",
            "# C reANGPTL3 induces FAK and JNK phosphorylation, and cilengitide can inhibit this phosphorylation.",
            "# D, E cilengitide or JNK inhibitor SP600125 can relieve the suppression of thermogenic genes by ANGPTL3.",
            "# F VO2 increases after injection of cilengitide in ob/ob mice.",
            "# G-I scWAT thermogenic gene and UCP1 protein are up-regulated, adipocytes become smaller, indicating enhanced browning.",
        ],
        "Figure_7": [
            "# A-C ATF4-activated Hepa1-6 conditioned medium inhibits adipocyte thermogenic genes and UCP1 protein.",
            "# D, E ATF4 is upregulated in the liver and VO2 is decreased after liver ATF4 plasmid injection.",
            "# F-H UCP1 protein and OCR decreased in scWAT, indicating that browning was inhibited.",
            "# I BAT maximum OCR decreases.",
            "# J, K ATF4 upregulation induces increased expression of CIRP and ANGPTL3, supporting the ATF4-CIRP-ANGPTL3 cascade.",
        ],
    }

    label_color = random.choice(LABEL_COLORS)
    for fig_num in range(1, 8):
        fig_key = f"Figure_{fig_num}"
        image_path = base_dir / f"{fig_key}.jpg"
        fig_title = figure_titles_zh.get(fig_key) or figure_titles.get(fig_key)
        fig_title = normalize_figure_title(fig_title) if fig_title else fig_title
        if fig_title:
            slide_title = f"{fig_key.replace('_', ' ')}: {fig_title}"
        else:
            slide_title = fig_key.replace("_", " ")
        fig_paragraphs = list(results[fig_key])
        fig_paragraphs.append(build_summary_line(fig_paragraphs))
        add_result_slide(
            prs,
            slide_title,
            image_path,
            fig_paragraphs,
            label_color,
            animator,
            transition_name,
            anim_presets,
        )

    strengths = [
        "Proposing the ATF4-CIRP-ANGPTL3 cascade as a new pathway to inhibit fat browning",
        "Verify causality through liver-specific manipulation and multi-level evidence chain",
        "Pointing to ANGPTL3/integrin αvβ3 as a potential intervention target",
    ]
    limits = [
        "Mainly based on mouse models, human evidence is still limited",
        "The upstream triggering factors of the molecular mechanism still need to be further verified.",
        "Insufficient assessment of metabolic safety of long-term interventions",
    ]
    add_summary_slide(prs, strengths, limits, label_color, animator, transition_name, anim_presets)

    prs.save(str(output_path))


if __name__ == "__main__":
    main()
